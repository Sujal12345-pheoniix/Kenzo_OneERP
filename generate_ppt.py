import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation and set widescreen (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme (Kenzo OneERP Brand Colors)
C_DARK_BG       = RGBColor(15, 23, 42)      # Slate 900
C_LIGHT_BG      = RGBColor(248, 250, 252)   # Slate 50
C_PRIMARY       = RGBColor(2, 132, 199)     # Sky 600
C_ACCENT        = RGBColor(99, 102, 241)    # Indigo 500
C_SUCCESS       = RGBColor(16, 185, 129)    # Emerald 500
C_TEXT_DARK     = RGBColor(15, 23, 42)      # Slate 900
C_TEXT_LIGHT    = RGBColor(255, 255, 255)   # White
C_TEXT_MUTED    = RGBColor(100, 116, 139)   # Slate 500
C_WHITE         = RGBColor(255, 255, 255)

FONT_HEADING = "Segoe UI"
FONT_BODY    = "Segoe UI"

def set_background(slide, color):
    """Sets a solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_text_box(slide, left, top, width, height, text="", font_name=FONT_BODY, font_size=Pt(14), font_color=C_TEXT_DARK, bold=False, align=PP_ALIGN.LEFT):
    """Utility to create a text box with custom styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return tf

def add_header(slide, title, category="KENZO ONEERP"):
    """Adds a standard header to a content slide."""
    # Category tracker
    tf_cat = create_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.3))
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = FONT_HEADING
    p_cat.font.size = Pt(10)
    p_cat.font.color.rgb = C_PRIMARY
    p_cat.font.bold = True
    
    # Title
    tf_title = create_text_box(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.6))
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(26)
    p_title.font.color.rgb = C_TEXT_DARK
    p_title.font.bold = True

    # Horizontal divider rule
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(226, 232, 240)
    shape.line.fill.background()

def add_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=None):
    """Draws a card background shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ==========================================
# SLIDE 1: Title & Overview (Dark Theme)
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
set_background(slide1, C_DARK_BG)

# Decorative vertical colored bar
accent_bar = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(0), Inches(0), Inches(0.4), Inches(7.5)
)
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = C_PRIMARY
accent_bar.line.fill.background()

# Subtitle / Tracker
tf_sub = create_text_box(slide1, Inches(1.2), Inches(2.2), Inches(10), Inches(0.4))
p_sub = tf_sub.paragraphs[0]
p_sub.text = "NEXT-GENERATION ENTERPRISE SUITE"
p_sub.font.name = FONT_HEADING
p_sub.font.size = Pt(12)
p_sub.font.color.rgb = C_PRIMARY
p_sub.font.bold = True

# Main Title
tf_main = create_text_box(slide1, Inches(1.2), Inches(2.6), Inches(11), Inches(1.8))
p_main = tf_main.paragraphs[0]
p_main.text = "Kenzo OneERP"
p_main.font.name = FONT_HEADING
p_main.font.size = Pt(54)
p_main.font.color.rgb = C_TEXT_LIGHT
p_main.font.bold = True

# System Upgrade Tagline
p_main2 = tf_main.add_paragraph()
p_main2.text = "Advanced Employee Portals & Smart Attendance Systems"
p_main2.font.name = FONT_HEADING
p_main2.font.size = Pt(22)
p_main2.font.color.rgb = C_ACCENT
p_main2.font.bold = False

# Description/Overview Footer
tf_desc = create_text_box(slide1, Inches(1.2), Inches(4.8), Inches(10), Inches(1.2))
p_desc = tf_desc.paragraphs[0]
p_desc.text = "An optimized enterprise resource system upgrading operational transparency. The new module introduces detailed interactive profiles for administrators to track progress, paired with a cryptographically enforced role-based check-in calendar limiting employee attendance to the present day."
p_desc.font.name = FONT_BODY
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = RGBColor(148, 163, 184) # slate 400

# ==========================================
# SLIDE 2: Problem Statement & Issues
# ==========================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide2, C_LIGHT_BG)
add_header(slide2, "Current Challenges & Issues")

# Left Column (The Pain Point)
left_card = add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), C_WHITE)
tf_left = create_text_box(slide2, Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
p_l_title = tf_left.paragraphs[0]
p_l_title.text = "The Administrative Gap"
p_l_title.font.name = FONT_HEADING
p_l_title.font.size = Pt(20)
p_l_title.font.color.rgb = C_PRIMARY
p_l_title.font.bold = True

p_l_body = tf_left.add_paragraph()
p_l_body.text = "\nModern enterprises require high levels of auditability. Managing hybrid teams without central task summaries, leaves track-sheets, and verified timestamps results in operational disconnect."
p_l_body.font.name = FONT_BODY
p_l_body.font.size = Pt(14)
p_l_body.font.color.rgb = C_TEXT_MUTED

# Right Column (Issues List)
right_card = add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), C_WHITE)
tf_right = create_text_box(slide2, Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))
p_r_title = tf_right.paragraphs[0]
p_r_title.text = "Key Operational Leaks"
p_r_title.font.name = FONT_HEADING
p_r_title.font.size = Pt(20)
p_r_title.font.color.rgb = C_TEXT_DARK
p_r_title.font.bold = True

issues = [
    ("Siloed Employee Data", "Management has to jump across departments, payrolls, and task boards to inspect a single team member's active output."),
    ("Attendance Back-dating & Manipulation", "Traditional timesheets allow users to check in retrospectively for past dates or prepopulate future days, leading to attendance leaks."),
    ("Admin Overheads & CEO Noise", "Administrators spend hours aggregating metrics, while executives (CEOs/Admins) are cluttered with daily check-in tasks they do not need.")
]

for title, desc in issues:
    p_t = tf_right.add_paragraph()
    p_t.text = f"\n• {title}"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_ACCENT
    
    p_d = tf_right.add_paragraph()
    p_d.text = desc
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = C_TEXT_MUTED

# ==========================================
# SLIDE 3: Why Needed
# ==========================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide3, C_LIGHT_BG)
add_header(slide3, "Strategic Rationale")

reasons = [
    ("Operational Auditability", "Guarantees a clean, unmanipulated paper trail of attendance records directly associated with project delivery timelines.", C_PRIMARY),
    ("Resource Synchronization", "Links employee tasks, leaves, and payroll directly to their profiles so managers can allocate resources in real-time.", C_ACCENT),
    ("Frictionless Compliance", "Automates policies like weekend lockouts, check-in thresholds, and executive role exemptions automatically without manual oversight.", C_SUCCESS)
]

for idx, (title, desc, color) in enumerate(reasons):
    left_pos = Inches(0.8 + idx * 3.9)
    # Card background
    add_card(slide3, left_pos, Inches(1.8), Inches(3.7), Inches(4.8), C_WHITE)
    
    # Border accent line
    top_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(3.7), Inches(0.1))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = color
    top_line.line.fill.background()
    
    tf_card = create_text_box(slide3, left_pos + Inches(0.3), Inches(2.2), Inches(3.1), Inches(4.0))
    
    p_title = tf_card.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = color
    
    p_desc = tf_card.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = C_TEXT_MUTED

# ==========================================
# SLIDE 4: Mention Features
# ==========================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide4, C_LIGHT_BG)
add_header(slide4, "Core Module Upgrades")

# Left: Admin Feature Set
add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), C_WHITE)
tf_feat_l = create_text_box(slide4, Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))

p_feat_l_title = tf_feat_l.paragraphs[0]
p_feat_l_title.text = "Admin Profile Control Hub"
p_feat_l_title.font.name = FONT_HEADING
p_feat_l_title.font.size = Pt(20)
p_feat_l_title.font.color.rgb = C_PRIMARY
p_feat_l_title.font.bold = True

admin_feats = [
    ("Unified Profile Dashboard", "A single-view card listing project details, tasks, attendance history grids, and leave records."),
    ("Task Progress Tracking", "Sorts and reviews pending vs completed tasks grouped by priority (Low, Medium, High, Critical)."),
    ("Historical Attendance Grid", "Inspects a interactive 30-day grid with status legends (Present, Late, Absent, Leave).")
]

for t, d in admin_feats:
    p_t = tf_feat_l.add_paragraph()
    p_t.text = f"\n✓ {t}"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_TEXT_DARK
    
    p_d = tf_feat_l.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = C_TEXT_MUTED

# Right: Employee Feature Set
add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), C_WHITE)
tf_feat_r = create_text_box(slide4, Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))

p_feat_r_title = tf_feat_r.paragraphs[0]
p_feat_r_title.text = "Strict Check-In Calendar"
p_feat_r_title.font.name = FONT_HEADING
p_feat_r_title.font.size = Pt(20)
p_feat_r_title.font.color.rgb = C_ACCENT
p_feat_r_title.font.bold = True

emp_feats = [
    ("Date-Restricted check-ins", "Employees can only check in for the current calendar date. Past and future checks are completely locked."),
    ("Exemption Policies", "Admin and CEO accounts are automatically exempted from daily checking in, reducing clutter."),
    ("Integrated Month Stats", "Real-time calendar updates with color-coded status tiles (Green=Present, Yellow=Late, Red=Absent, Violet=Leave).")
]

for t, d in emp_feats:
    p_t = tf_feat_r.add_paragraph()
    p_t.text = f"\n✓ {t}"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = C_TEXT_DARK
    
    p_d = tf_feat_r.add_paragraph()
    p_d.text = d
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(11)
    p_d.font.color.rgb = C_TEXT_MUTED


# ==========================================
# SLIDE 5: Stack & Tools
# ==========================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide5, C_LIGHT_BG)
add_header(slide5, "Technology Stack & Architecture")

techs = [
    ("Next.js App Router", "React-based web framework optimizing client rendering and server actions.", "Frontend / Routing"),
    ("Prisma ORM", "Type-safe database client mapping PostgreSQL data relations seamlessly.", "Data Layer"),
    ("PostgreSQL", "Robust SQL database enforcing tenant isolation and referential integrity.", "Database"),
    ("Lucide Icons", "Vector library providing clean, uniform indicators for active modules.", "Visuals"),
    ("Bcrypt.js / JWT", "Cryptographic token verification ensuring roles and session safety.", "Security"),
    ("TailwindCSS", "Utility styling framework facilitating standard brand themes and responsiveness.", "Design System")
]

for idx, (title, desc, layer) in enumerate(techs):
    col = idx % 3
    row = idx // 3
    
    l_pos = Inches(0.8 + col * 3.9)
    t_pos = Inches(1.8 + row * 2.5)
    
    add_card(slide5, l_pos, t_pos, Inches(3.7), Inches(2.2), C_WHITE)
    
    tf_tech = create_text_box(slide5, l_pos + Inches(0.3), t_pos + Inches(0.2), Inches(3.1), Inches(1.8))
    
    p_layer = tf_tech.paragraphs[0]
    p_layer.text = layer.upper()
    p_layer.font.name = FONT_HEADING
    p_layer.font.size = Pt(9)
    p_layer.font.bold = True
    p_layer.font.color.rgb = C_PRIMARY
    
    p_title = tf_tech.add_paragraph()
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = C_TEXT_DARK
    
    p_desc = tf_tech.add_paragraph()
    p_desc.text = desc
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(11)
    p_desc.font.color.rgb = C_TEXT_MUTED

# ==========================================
# SLIDE 6: Workflow & Data flow Diagram
# ==========================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide6, C_LIGHT_BG)
add_header(slide6, "System Flow & Verification Model")

# Visual Process Blocks (Horizontal Flow)
blocks = [
    ("1. Request Session", "Validates JWT & reads user role", C_PRIMARY),
    ("2. Role Verification", "Is user Admin/CEO?", C_ACCENT),
    ("3. Date Boundary", "Compare checkIn date with Today", C_PRIMARY),
    ("4. Database Lock", "Create unique Attendance entry", C_SUCCESS)
]

for idx, (step, desc, color) in enumerate(blocks):
    left_pos = Inches(0.8 + idx * 2.9)
    # Block Shape
    block = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.5), Inches(2.6), Inches(3.0))
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()
    
    tf_b = create_text_box(slide6, left_pos + Inches(0.2), Inches(2.7), Inches(2.2), Inches(2.6))
    
    p_step = tf_b.paragraphs[0]
    p_step.text = step
    p_step.font.name = FONT_HEADING
    p_step.font.size = Pt(14)
    p_step.font.bold = True
    p_step.font.color.rgb = C_TEXT_LIGHT
    
    p_desc = tf_b.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = RGBColor(241, 245, 249) # Light slate
    
    # Arrow Connector (if not last)
    if idx < 3:
        arrow = slide6.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, 
            left_pos + Inches(2.65), Inches(3.7), Inches(0.2), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(148, 163, 184)
        arrow.line.fill.background()

# Technical Detail Banner
add_card(slide6, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.8), C_DARK_BG)
tf_banner = create_text_box(slide6, Inches(1.1), Inches(6.15), Inches(11.1), Inches(0.5))
p_banner = tf_banner.paragraphs[0]
p_banner.text = "Strict Constraint: Next.js API route /api/attendance rejects check-in timestamps not matching the server's current date (UTC normalized to local time zone) with HTTP status 400."
p_banner.font.name = FONT_BODY
p_banner.font.size = Pt(11)
p_banner.font.color.rgb = RGBColor(226, 232, 240)
p_banner.font.bold = True

# ==========================================
# SLIDE 7: Advantages
# ==========================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide7, C_LIGHT_BG)
add_header(slide7, "Operational Advantages")

advs = [
    ("Anti-Tamper Design", "Secures operations by removing backdated timesheets and pre-filled sheets.", C_PRIMARY),
    ("Executive Clarity", "Eliminates visual noise for Admins/CEOs by excluding them from daily check-in workflows.", C_ACCENT),
    ("Real-Time Analytics", "Syncs attendance logs with ongoing tasks for accurate operational productivity charts.", C_SUCCESS),
    ("Scalable Structure", "Prisma schema ensures smooth data relations, handling future custom modules.", C_PRIMARY)
]

for idx, (title, desc, color) in enumerate(advs):
    col = idx % 2
    row = idx // 2
    
    l_pos = Inches(0.8 + col * 5.9)
    t_pos = Inches(1.8 + row * 2.5)
    
    add_card(slide7, l_pos, t_pos, Inches(5.6), Inches(2.2), C_WHITE)
    
    # Left vertical indicator bar
    indicator = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, l_pos, t_pos, Inches(0.1), Inches(2.2))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    
    tf_adv = create_text_box(slide7, l_pos + Inches(0.4), t_pos + Inches(0.3), Inches(4.8), Inches(1.6))
    
    p_title = tf_adv.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = C_TEXT_DARK
    
    p_desc = tf_adv.add_paragraph()
    p_desc.text = f"\n{desc}"
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = C_TEXT_MUTED

# ==========================================
# SLIDE 8: Demos of Project
# ==========================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide8, C_LIGHT_BG)
add_header(slide8, "Visual Interface Demos")

# Demo 1 Card
add_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), C_WHITE)
tf_demo1 = create_text_box(slide8, Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))

p_d1_title = tf_demo1.paragraphs[0]
p_d1_title.text = "1. Admin Profile Portal"
p_d1_title.font.name = FONT_HEADING
p_d1_title.font.size = Pt(20)
p_d1_title.font.color.rgb = C_PRIMARY
p_d1_title.font.bold = True

d1_bullets = [
    "Interactive Avatars: Clickable member bubbles (e.g., SK, AS, CS) in Organization view.",
    "Integrated Metrics: Displays active project list, priority-sorted tasks, and leave logs.",
    "Grid Timeline: A colorful visual grid mapping the last 30 days of attendance entries."
]
for bullet in d1_bullets:
    p = tf_demo1.add_paragraph()
    p.text = f"\n• {bullet}"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = C_TEXT_MUTED

# Demo 2 Card
add_card(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), C_WHITE)
tf_demo2 = create_text_box(slide8, Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))

p_d2_title = tf_demo2.paragraphs[0]
p_d2_title.text = "2. Smart Attendance Calendar"
p_d2_title.font.name = FONT_HEADING
p_d2_title.font.size = Pt(20)
p_d2_title.font.color.rgb = C_ACCENT
p_d2_title.font.bold = True

d2_bullets = [
    "Single-Day Access: Present day cell highlights in active blue, prompting check-in.",
    "Calendar Lockout: Yesterday's cells and future dates are grayed-out and disabled.",
    "Real-time Indicators: Color keys update immediately upon action (Green for Present, Yellow for Late)."
]
for bullet in d2_bullets:
    p = tf_demo2.add_paragraph()
    p.text = f"\n• {bullet}"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = C_TEXT_MUTED

# ==========================================
# SLIDE 9: Thanks & Feedback (Dark Theme)
# ==========================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide9, C_DARK_BG)

tf_thanks = create_text_box(slide9, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.0), align=PP_ALIGN.CENTER)

p_t1 = tf_thanks.paragraphs[0]
p_t1.text = "Thank You!"
p_t1.font.name = FONT_HEADING
p_t1.font.size = Pt(54)
p_t1.font.color.rgb = C_TEXT_LIGHT
p_t1.font.bold = True
p_t1.alignment = PP_ALIGN.CENTER

p_t2 = tf_thanks.add_paragraph()
p_t2.text = "Q&A and Feedback Session"
p_t2.font.name = FONT_HEADING
p_t2.font.size = Pt(24)
p_t2.font.color.rgb = C_PRIMARY
p_t2.font.bold = True
p_t2.alignment = PP_ALIGN.CENTER

p_t3 = tf_thanks.add_paragraph()
p_t3.text = "\n\nWe welcome your feedback on the new Kenzo OneERP features."
p_t3.font.name = FONT_BODY
p_t3.font.size = Pt(16)
p_t3.font.color.rgb = RGBColor(148, 163, 184)
p_t3.alignment = PP_ALIGN.CENTER

# Save the presentation
output_path = "Kenzo_OneERP_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {os.path.abspath(output_path)}")
